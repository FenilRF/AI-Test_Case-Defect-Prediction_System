"""
Document Extractor Service
------------------------------
Extracts text and UI-relevant structure from uploaded documents:
  • PDF  → PyPDF2
  • DOCX → python-docx
  • PPTX → python-pptx
  • TXT  → direct read

After extraction:
  - Identifies UI keywords (field, button, submit, login, etc.)
  - Detects module boundaries from headings / sections
  - Detects flows from sequential keyword patterns
"""

import logging
import os
import re
from typing import Any, Dict, List

logger = logging.getLogger(__name__)

# ── UI keyword bank ──────────────────────────────────────────
_UI_KEYWORDS = [
    "form", "button", "input", "field", "textbox", "textarea",
    "dropdown", "select", "checkbox", "radio", "toggle", "switch",
    "slider", "submit", "cancel", "save", "delete", "edit",
    "login", "logout", "register", "signup", "sign up", "sign in",
    "dashboard", "sidebar", "navbar", "navigation", "menu",
    "header", "footer", "modal", "dialog", "popup", "toast",
    "alert", "notification", "error", "warning", "success",
    "table", "grid", "list", "card", "panel", "tab", "tabs",
    "breadcrumb", "pagination", "search", "filter", "sort",
    "upload", "download", "file", "image", "avatar", "icon",
    "tooltip", "placeholder", "label", "link", "anchor",
    "date picker", "calendar", "stepper", "progress bar",
    "spinner", "loading", "skeleton", "badge", "chip", "tag",
]

# ── Module boundary markers ──────────────────────────────────
_MODULE_MARKERS = [
    "module", "page", "screen", "section", "component",
    "feature", "view", "panel", "form", "wizard", "flow",
]

# ── Flow keyword patterns ────────────────────────────────────
_FLOW_KEYWORDS = [
    "login", "register", "signup", "dashboard", "home",
    "create", "add", "edit", "update", "delete", "remove",
    "search", "filter", "export", "import", "download", "upload",
    "submit", "approve", "reject", "cancel", "confirm",
    "checkout", "payment", "cart", "order",
    "profile", "settings", "account", "password",
]


def _extract_text_from_pdf(file_path: str) -> List[str]:
    """Extract text from each page of a PDF."""
    pages = []
    try:
        import PyPDF2
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text and text.strip():
                    pages.append(text.strip())
    except ImportError:
        logger.warning("PyPDF2 not installed, cannot extract PDF text")
    except Exception as e:
        logger.warning("PDF extraction failed for %s: %s", file_path, e)
    return pages


def _extract_text_from_docx(file_path: str) -> List[str]:
    """Extract text from a DOCX file."""
    pages = []
    try:
        import docx
        doc = docx.Document(file_path)
        current_section = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                if current_section:
                    pages.append("\n".join(current_section))
                    current_section = []
                continue
            # Check if it's a heading (potential module boundary)
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                if current_section:
                    pages.append("\n".join(current_section))
                    current_section = []
            current_section.append(text)
        if current_section:
            pages.append("\n".join(current_section))

        # Also extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    table_text.append(" | ".join(row_text))
            if table_text:
                pages.append("Table:\n" + "\n".join(table_text))
    except ImportError:
        logger.warning("python-docx not installed, cannot extract DOCX text")
    except Exception as e:
        logger.warning("DOCX extraction failed for %s: %s", file_path, e)
    return pages


def _extract_text_from_pptx(file_path: str) -> List[str]:
    """Extract text from each slide of a PPTX."""
    pages = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                pages.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
    except ImportError:
        logger.warning("python-pptx not installed, cannot extract PPTX text")
    except Exception as e:
        logger.warning("PPTX extraction failed for %s: %s", file_path, e)
    return pages


def _extract_text_from_txt(file_path: str) -> List[str]:
    """Read text from a TXT file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read().strip()
        if content:
            # Split into sections by double newlines or heading patterns
            sections = re.split(r"\n{2,}|\n(?=[A-Z][A-Za-z\s]+:)", content)
            return [s.strip() for s in sections if s.strip()]
    except Exception as e:
        logger.warning("TXT read failed for %s: %s", file_path, e)
    return []


def _detect_ui_keywords(text: str) -> List[str]:
    """Find UI-relevant keywords in text."""
    text_lower = text.lower()
    found = []
    for kw in _UI_KEYWORDS:
        pattern = rf"\b{re.escape(kw)}\b"
        if re.search(pattern, text_lower):
            found.append(kw)
    return sorted(set(found))


def _detect_modules(pages_text: List[str]) -> List[str]:
    """Detect module/section boundaries from page text."""
    modules = set()
    for page in pages_text:
        text_lower = page.lower()
        # Look for heading-like patterns
        headings = re.findall(
            r"(?:^|\n)\s*(?:\d+\.?\s+)?([A-Z][A-Za-z\s]+?)(?:\s*[:—\-\n])",
            page,
        )
        for h in headings[:5]:
            h = h.strip()
            if 3 < len(h) < 60:
                modules.add(h)

        # Also check for module marker keywords
        for marker in _MODULE_MARKERS:
            pattern = rf"\b{re.escape(marker)}\s*[:—\-]?\s*([A-Za-z\s]{{3,40}})"
            matches = re.findall(pattern, text_lower)
            for m in matches[:3]:
                modules.add(m.strip().title())

    return sorted(modules)[:20]


def _detect_flows(ui_keywords: List[str]) -> List[Dict[str, Any]]:
    """Detect potential user flows from keyword patterns."""
    flows = []
    # Check for common flow patterns
    flow_patterns = [
        {
            "name": "Authentication Flow",
            "keywords": ["login", "register", "signup", "password", "sign in", "sign up"],
            "steps": ["Register/Sign Up", "Login/Sign In", "Dashboard"],
        },
        {
            "name": "CRUD Flow",
            "keywords": ["create", "add", "edit", "update", "delete", "remove"],
            "steps": ["Create/Add", "View/List", "Edit/Update", "Delete/Remove"],
        },
        {
            "name": "Search & Filter Flow",
            "keywords": ["search", "filter", "sort", "export"],
            "steps": ["Search", "Apply Filters", "Sort Results", "Export Data"],
        },
        {
            "name": "Form Submission Flow",
            "keywords": ["form", "input", "submit", "save", "cancel"],
            "steps": ["Fill Form", "Validate", "Submit", "Confirmation"],
        },
        {
            "name": "File Management Flow",
            "keywords": ["upload", "download", "file", "image"],
            "steps": ["Select File", "Upload", "Preview", "Download"],
        },
        {
            "name": "E-Commerce Flow",
            "keywords": ["cart", "checkout", "payment", "order"],
            "steps": ["Browse Products", "Add to Cart", "Checkout", "Payment", "Order Confirmation"],
        },
    ]

    kw_set = set(ui_keywords)
    for pattern in flow_patterns:
        overlap = kw_set.intersection(set(pattern["keywords"]))
        if len(overlap) >= 2:
            flows.append({
                "name": pattern["name"],
                "matched_keywords": sorted(overlap),
                "steps": pattern["steps"],
            })

    return flows


def extract_from_document(file_path: str) -> Dict[str, Any]:
    """
    Master document extraction function.

    Parameters
    ----------
    file_path : str
        Absolute path to the uploaded document.

    Returns
    -------
    dict
        {
            "pages_text": [...],
            "total_pages": int,
            "ui_keywords": [...],
            "modules_detected": [...],
            "flows": [...],
            "file_type": str,
        }
    """
    ext = os.path.splitext(file_path)[1].lower()
    logger.info("Extracting text from %s (type: %s)", file_path, ext)

    # Route to appropriate extractor
    extractors = {
        ".pdf": _extract_text_from_pdf,
        ".docx": _extract_text_from_docx,
        ".doc": _extract_text_from_docx,
        ".pptx": _extract_text_from_pptx,
        ".txt": _extract_text_from_txt,
    }

    extractor = extractors.get(ext)
    if not extractor:
        logger.warning("Unsupported document type: %s", ext)
        return {
            "pages_text": [],
            "total_pages": 0,
            "ui_keywords": [],
            "modules_detected": [],
            "flows": [],
            "file_type": ext,
        }

    pages_text = extractor(file_path)
    combined_text = "\n\n".join(pages_text)

    # Analyze extracted text
    ui_keywords = _detect_ui_keywords(combined_text)
    modules = _detect_modules(pages_text)
    flows = _detect_flows(ui_keywords)

    logger.info(
        "Document extraction complete: pages=%d, keywords=%d, modules=%d, flows=%d",
        len(pages_text), len(ui_keywords), len(modules), len(flows),
    )

    return {
        "pages_text": pages_text,
        "total_pages": len(pages_text),
        "ui_keywords": ui_keywords,
        "modules_detected": modules,
        "flows": flows,
        "file_type": ext,
    }
